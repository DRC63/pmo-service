"""P3MAI-branded python-pptx helper for the Method Map summary decks. 16:9."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

NAVY = RGBColor(0x0B, 0x25, 0x45)
NAVYL = RGBColor(0x1B, 0x3F, 0x6E)
GOLD = RGBColor(0xC9, 0xA2, 0x27)
GOLDD = RGBColor(0xA8, 0x84, 0x1C)
GREEN = RGBColor(0x2E, 0x7D, 0x5B)
RED = RGBColor(0xC0, 0x39, 0x2B)
PURPLE = RGBColor(0x8E, 0x5B, 0xE0)
GREY = RGBColor(0x5B, 0x66, 0x75)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
TEXT = RGBColor(0x1C, 0x2B, 0x3A)
BG = RGBColor(0xF6, 0xF7, 0xF9)
HEAD_FONT = "Segoe UI Semibold"
BODY_FONT = "Segoe UI"

EMU_W = Inches(13.333)
EMU_H = Inches(7.5)


class Deck:
    def __init__(self, doc_id, classification="OFFICIAL"):
        self.prs = Presentation()
        self.prs.slide_width = EMU_W
        self.prs.slide_height = EMU_H
        self.doc_id = doc_id
        self.classification = classification
        self.n = 0

    def _blank(self):
        return self.prs.slides.add_slide(self.prs.slide_layouts[6])

    def _rect(self, s, x, y, w, h, fill, line=None):
        sp = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
        sp.fill.solid(); sp.fill.fore_color.rgb = fill
        if line is None:
            sp.line.fill.background()
        else:
            sp.line.color.rgb = line; sp.line.width = Pt(1)
        sp.shadow.inherit = False
        return sp

    def _text(self, s, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, sp_after=6):
        tb = s.shapes.add_textbox(x, y, w, h); tf = tb.text_frame
        tf.word_wrap = True; tf.vertical_anchor = anchor
        first = True
        for item in runs:
            txt, size, color, bold, font = (item + (None,))[:5]
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            p.alignment = align; p.space_after = Pt(sp_after); p.space_before = Pt(0)
            r = p.add_run(); r.text = txt
            r.font.size = Pt(size); r.font.bold = bold
            r.font.color.rgb = color; r.font.name = font or BODY_FONT
        return tb

    def _footer(self, s):
        self._text(s, Inches(0.35), Inches(7.02), Inches(4), Inches(0.4),
                   [(self.classification, 9, GOLDD, True)])
        self._text(s, Inches(11.8), Inches(7.02), Inches(1.2), Inches(0.4),
                   [(f"{self.doc_id}  ·  {self.n}", 9, GREY, False)], align=PP_ALIGN.RIGHT)

    def _titlebar(self, s, title):
        # slim navy bar + gold rule, title in navy below
        self._text(s, Inches(0.6), Inches(0.45), Inches(12), Inches(1.0),
                   [(title, 26, NAVY, True, HEAD_FONT)])
        rule = self._rect(s, Inches(0.62), Inches(1.32), Inches(3.2), Pt(3), GOLD)
        return rule

    def title_slide(self, title, subtitle):
        s = self._blank(); self.n += 1
        self._rect(s, 0, 0, EMU_W, EMU_H, NAVY)
        self._rect(s, 0, Inches(4.55), EMU_W, Pt(3), GOLD)
        self._text(s, Inches(0.9), Inches(1.5), Inches(11.5), Inches(0.7),
                   [("P3MAI  ·  METHOD MAP", 16, GOLD, True, HEAD_FONT)])
        self._text(s, Inches(0.9), Inches(2.6), Inches(11.5), Inches(1.8),
                   [(title, 40, WHITE, True, HEAD_FONT)])
        self._text(s, Inches(0.9), Inches(4.75), Inches(11.5), Inches(0.8),
                   [(subtitle, 18, RGBColor(0xC8, 0xD2, 0xE0), False)])
        self._text(s, Inches(0.9), Inches(6.4), Inches(11.5), Inches(0.6),
                   [(f"{self.doc_id}  ·  Summary  ·  v1.0  ·  1 August 2026  ·  {self.classification}",
                     12, RGBColor(0x9F, 0xB0, 0xC4), False)])
        return s

    def bullets(self, title, items, lead=None):
        s = self._blank(); self.n += 1
        self._titlebar(s, title)
        y = 1.75
        if lead:
            self._text(s, Inches(0.62), Inches(y), Inches(12.1), Inches(0.7),
                       [(lead, 15, GREY, False)])
            y += 0.85
        runs = []
        for it in items:
            if isinstance(it, tuple):
                txt, color = it
            else:
                txt, color = it, TEXT
            runs.append(("•  " + txt, 17, color, False))
        self._text(s, Inches(0.7), Inches(y), Inches(12.0), Inches(7.0 - y - 0.6), runs, sp_after=12)
        self._footer(s)
        return s

    def image(self, title, path, lead=None):
        s = self._blank(); self.n += 1
        self._titlebar(s, title)
        top = Inches(1.85)
        if lead:
            self._text(s, Inches(0.62), Inches(1.7), Inches(12.1), Inches(0.6),
                       [(lead, 14, GREY, False)])
            top = Inches(2.25)
        # size image to fit width, centred
        from PIL import Image
        iw, ih = Image.open(path).size
        max_w = Inches(11.6); max_h = Inches(4.7)
        ratio = min(max_w / iw, max_h / ih)
        w = Emu(int(iw * ratio)); h = Emu(int(ih * ratio))
        left = Emu(int((EMU_W - w) / 2))
        s.shapes.add_picture(path, left, top, width=w, height=h)
        self._footer(s)
        return s

    def table(self, title, headers, rows, col_widths=None, lead=None):
        s = self._blank(); self.n += 1
        self._titlebar(s, title)
        y = Inches(1.85)
        if lead:
            self._text(s, Inches(0.62), Inches(1.7), Inches(12.1), Inches(0.6),
                       [(lead, 14, GREY, False)])
            y = Inches(2.3)
        nrows = len(rows) + 1; ncols = len(headers)
        total_w = Inches(12.1)
        height = Inches(min(4.6, 0.5 + 0.5 * len(rows)))
        gtbl = s.shapes.add_table(nrows, ncols, Inches(0.62), y, total_w, height).table
        if col_widths:
            for i, cw in enumerate(col_widths):
                gtbl.columns[i].width = Inches(cw)
        for j, h in enumerate(headers):
            c = gtbl.cell(0, j); c.fill.solid(); c.fill.fore_color.rgb = NAVY
            tf = c.text_frame; tf.word_wrap = True; p = tf.paragraphs[0]
            r = p.add_run(); r.text = h; r.font.bold = True; r.font.size = Pt(12)
            r.font.color.rgb = WHITE; r.font.name = HEAD_FONT
        for ri, row in enumerate(rows, start=1):
            for j, val in enumerate(row):
                c = gtbl.cell(ri, j)
                c.fill.solid(); c.fill.fore_color.rgb = WHITE if ri % 2 else BG
                tf = c.text_frame; tf.word_wrap = True; p = tf.paragraphs[0]
                r = p.add_run(); r.text = str(val); r.font.size = Pt(11)
                r.font.color.rgb = TEXT; r.font.name = BODY_FONT
        self._footer(s)
        return s

    def save(self, path):
        self.prs.save(path)
